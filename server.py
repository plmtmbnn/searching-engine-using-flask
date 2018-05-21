from flask import Flask, render_template, flash, request, send_file
from wtforms import Form, TextField, TextAreaField, validators, StringField, SubmitField
from flask_bootstrap import Bootstrap
from collections import defaultdict
import math
import sys
import os
import time
import PyPDF2
from docx import Document
from pptx import Presentation
from operator import itemgetter

DEBUG = True
app = Flask(__name__)
app.config.from_object(__name__)
app.config['SECRET_KEY'] = '7d441f27d441f27567d441f2b6176a'
Bootstrap(app)

list_document = []
documents = {}
total_document_found = 0
file_type = [".pdf",".txt",".docx",".pptx"]
document_filenames = {}
document_paths = {}
N = 0
dictionary = set()
postings = defaultdict(dict)
document_frequency = defaultdict(int)
length = defaultdict(float)
characters = " .,!#$%^&*();:\n\t\\\"?!{}[]<>-0987654321"
is_main_function_had_called = False

def main():
    global is_main_function_had_called
    if not is_main_function_had_called:
        start_time = time.time()
        read_driver("E:\\")
        initialize_terms_and_postings()
        initialize_document_frequencies()
        initialize_lengths()
        is_main_function_had_called = True
        print N
        print("\n--- %s seconds ---" % (time.time() - start_time))

def getDocumentContent(path):
    content = ""
    if(path.rfind(".pdf") == len(path) - 4):
        p = file(path, "rb")
        pdf_content = PyPDF2.PdfFileReader(p)
        for i in range(0, pdf_content.numPages):
            content += pdf_content.getPage(i).extractText() + "\n"
        content = " ".join(content.replace(u"\xa0", " ").strip().split())
    if(path.rfind(".txt") == len(path) - 4):
        f = open(path,'r')
        content = f.read()
        f.close()
    if(path.rfind(".docx") == len(path) - 5):
        doc = Document(path)
        for paragraph in doc.paragraphs:
            content += paragraph.text + "\n"
    if(path.rfind(".pptx") == len(path) - 5):
        pptx = Presentation(path)
        for slide in pptx.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                content += shape.text_frame.text + "\n"
    return content

def read_driver(source_root):
    id = len(document_filenames)
    global N, file_type
    for extension in file_type:
        for (dirpath, dirnames, filenames) in os.walk(source_root):
            for filename in filenames:
                if os.path.splitext(filename)[-1] == extension:
                    string =  dirpath+"/"+filename
                    document_paths[id] =filename
                    document_filenames[id] = string.replace("\\","/")
                    id += 1
        N = id + 1

def tokenize(document):
    """Returns a list whose elements are the separate terms in
    document.  Something of a hack, but for the simple documents we're
    using, it's okay.  Note that we case-fold when we tokenize, i.e.,
    we lowercase everything."""
    terms = document.lower().split()
    space = ' '
    return [term.strip(characters) for term in terms if term not in space]

def initialize_terms_and_postings():
    """Reads in each document in document_filenames, splits it into a
    list of terms (i.e., tokenizes it), adds new terms to the global
    dictionary, and adds the document to the posting list for each
    term, with value equal to the frequency of the term in the
    document."""
    global dictionary, postings
    for id in document_filenames:
        document = getDocumentContent(document_filenames[id])
        if(document_filenames[id].rfind(".pdf") == len(document_filenames[id]) - 4):
            terms = tokenize(document.encode('utf-8'))
        if(document_filenames[id].rfind(".txt") == len(document_filenames[id]) - 4):
            terms = tokenize(document)
        if(document_filenames[id].rfind(".docx") == len(document_filenames[id]) - 5):
            terms = tokenize(document)
        if(document_filenames[id].rfind(".pptx") == len(document_filenames[id]) - 5):
            terms = tokenize(document)
        unique_terms = set(terms)
        dictionary = dictionary.union(unique_terms)
        for term in unique_terms:
            postings[term][id] = terms.count(term) # the value is the
                                                   # frequency of the
                                                   # term in the
                                                   # document

def inverse_document_frequency(term):
    """Returns the inverse document frequency of term.  Note that if
    term isn't in the dictionary then it returns 0, by convention."""
    if term in dictionary:
        return math.log(N/document_frequency[term],2)
    else:
        return 0.0

def imp(term,id):
    """Returns the importance of term in document id.  If the term
    isn't in the document, then return 0."""
    if id in postings[term]:
        return postings[term][id]*inverse_document_frequency(term)
    else:
        return 0.0

def initialize_document_frequencies():
    """For each term in the dictionary, count the number of documents
    it appears in, and store the value in document_frequncy[term]."""
    global document_frequency
    for term in dictionary:
        document_frequency[term] = len(postings[term])

def initialize_lengths():
    """Computes the length for each document."""
    global length
    for id in document_filenames:
        l = 0
        for term in dictionary:
            l += imp(term,id)**2
        length[id] = math.sqrt(l)

def do_search(queries):
    """Asks the user what they would like to search for, and returns a
    list of relevant documents, in decreasing order of cosine
    similarity."""
    global documents, error, list_document
    results = {}
    query = tokenize(queries)
    if query == []:
        sys.exit()
    # find document ids containing all query terms.  Works by
    # intersecting the posting lists for all query terms.
    relevant_document_ids = intersection(
            [set(postings[term].keys()) for term in query])
    if not relevant_document_ids:
        documents.clear()
        list_document[:] = []
        flash('empty')
    else:
        scores = sorted([(id,similarity(query,id))
                         for id in relevant_document_ids],
                        key=lambda x: x[1],
                        reverse=True)
        print "Score: filename"
        global total_document_found
        total_document_found = 0
        for (id,score) in scores:
            print str(score)+": "+document_filenames[id]
            results[document_filenames[id]] = score
            total_document_found += 1
        flash("Total document found : " + str(total_document_found) + " of " + str(N))
    return results

def intersection(sets):
    """Returns the intersection of all sets in the list sets. Requires
    that the list sets contains at least one element, otherwise it
    raises an error."""
    return reduce(set.intersection, [s for s in sets])

def similarity(query,id):
    """Returns the cosine similarity between query and document id.
    Note that we don't bother dividing by the length of the query
    vector, since this doesn't make any difference to the ordering of
    search results."""
    similarity = 0.0
    for term in query:
        if term in dictionary:
            similarity += inverse_document_frequency(term)*imp(term,id)
    similarity = similarity / length[id]
    return similarity

class ReusableForm(Form):
    query = TextField('query:', validators=[validators.required()])

@app.route('/', methods=['GET', 'POST'])
def index():
    global documents, error, list_document
    form = ReusableForm(request.form)
    print form.errors
    main()
    if request.method == 'POST':
        query=request.form['query']
        if form.validate():
            documents.clear()
            list_document[:] = []
            for key, value in sorted(do_search(query).items(), key = itemgetter(1), reverse=True):
                documents[key] = value
                list_document.append(key)
                print key," => ",value
        else:
            documents.clear()
            list_document[:] = []
            flash('Error: The search field is cannot be blank.')
    else:
         documents.clear()
         list_document[:] = []
    return render_template('index.html', form=form, documents = documents, list_document =list_document)

@app.route('/download/<path:filename>', methods=['GET'])
def download(filename):
    return send_file(filename,as_attachment=True)

if __name__ == '__main__':
    app.run()
    app.jinja_env.auto_reload = True
    app.config['TEMPLATES_AUTO_RELOAD'] = True
